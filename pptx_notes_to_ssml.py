# we need to import that
import collections 
import collections.abc

from pptx import Presentation

delimiter = '# slide'

breaks = {
    "#break 500ms": '<break time="500ms" />',
    "#break 1s": '<break time="1000ms" />',
    "#break 2s": '<break time="2000ms" />',
    "#break 4s": '<break time="4000ms" />',
    "#break 5s": '<break time="5000ms" />',
    "#break 8s": '<break time="8000ms" />'
}

def extract_notes(filename: str) -> list:
    ppt = Presentation(filename)
    notes = []
    for page, slide in enumerate(ppt.slides):
        # this is the notes that doesn't appear on the ppt slide,
        # but really the 'presenter' note. 
        if slide.has_notes_slide:
            text_note = slide.notes_slide.notes_text_frame.text
            notes.append((page, text_note))
    return notes


def break_into_subslides(notes: list) -> list:
    subslide_notes = []
    for page_notes in notes:
        _, text = page_notes
        if text:
            subslides = text.split(delimiter)
            subslide_notes.extend(subslides)
    return subslide_notes


VOICE_JENNY = 'en-US-JennyNeural'


def build_ssml(notes: list[str], voice: str = VOICE_JENNY) -> list:

    items = []
    for i, note in enumerate(notes):
        note_text = note.strip()
        for break_key, break_value in breaks.items():
            note_text = note_text.replace(break_key, break_value)
        items.append(note)
        items.append(f'<bookmark mark="slide_{i}"/>')

    items_xml = "\n".join(['  ' + x for x in items])

    return f"""<speak xmlns="http://www.w3.org/2001/10/synthesis" xmlns:mstts="http://www.w3.org/2001/mstts" xmlns:emo="http://www.w3.org/2009/10/emotionml" version="1.0" xml:lang="en-US"> 
    <voice name="{voice}">
        {items_xml}
    </voice>
</speak>"""


def pptx_to_ssml(pptx_file: str, voice: str = VOICE_JENNY) -> str:
    notes = extract_notes(pptx_file)
    print(notes)
    subslide_notes = break_into_subslides(notes)
    print(subslide_notes)
    ssml = build_ssml(subslide_notes, voice)
    print(ssml)
    return ssml


def text_to_ssml(text: str, voice: str = VOICE_JENNY) -> str:
    subslides = text.split(delimiter)
    print(subslides)
    ssml = build_ssml(subslides, voice)
    print(ssml)
    return ssml


def txt_file_to_ssml(txt_file: str, voice: str) -> str:
    with open(txt_file, 'r') as f:
        text = f.read()
    return text_to_ssml(text, voice)



if __name__ == '__main__':
    ssml = pptx_to_ssml('test_slides.pptx')
    with open('audio.xml', 'w') as f:
        f.write(ssml)
