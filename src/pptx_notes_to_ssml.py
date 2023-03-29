# we need to import that
import collections 
import collections.abc
import sys

from pptx import Presentation

delimiter = '# slide'

breaks = {
    "#break 500ms": '<break time="500ms" />',
    "#break 1s": '<break time="1000ms" />',
    "#break 2s": '<break time="2000ms" />',
    "#break 4s": '<break time="4000ms" />',
    "#break 5s": '<break time="5000ms" />',
    "#break 8s": '<break time="8000ms" />'
}

def extract_notes(filename: str) -> list[str]:
    ppt = Presentation(filename)
    notes = []
    for _, slide in enumerate(ppt.slides):
        # this is the notes that doesn't appear on the ppt slide,
        # but really the 'presenter' note. 
        if slide.has_notes_slide:
            text_note = slide.notes_slide.notes_text_frame.text
            notes.append(text_note)
    return notes

def pptx_to_notes_str(input_file_name: str) -> str:
    notes = extract_notes(input_file_name)
    return "\n# slide\n".join(notes)


if __name__ == '__main__':
    input_file_name = 'testdata/test_script.pptx'
    output_file_name = 'testdata/output_ppt_notes.txt'
    
    if len(sys.argv) > 1:
        input_file_name = sys.argv[1]
    if len(sys.argv) > 2:
        output_file_name = sys.argv[2]
    
    result_text = pptx_to_notes_str(input_file_name)
    print(result_text)
    with open(output_file_name, 'w') as f:
        f.write(result_text)
